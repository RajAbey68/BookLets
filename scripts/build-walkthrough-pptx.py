"""Build the BookLets walkthrough deck as a .pptx file.

Run:  python3 scripts/build-walkthrough-pptx.py
Output: docs/booklets-walkthrough.pptx

Mirrors the content of docs/booklets-walkthrough.html so the
bookkeeper/accountant pack lives in a format they actually open.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette mirrors the HTML deck.
INK        = RGBColor(0x0F, 0x17, 0x2A)
INK_SOFT   = RGBColor(0x47, 0x55, 0x69)
ACCENT     = RGBColor(0x25, 0x63, 0xEB)
ACCENT_BG  = RGBColor(0xDB, 0xEA, 0xFE)
RULE       = RGBColor(0xE2, 0xE8, 0xF0)
CARD_BG    = RGBColor(0xF8, 0xFA, 0xFC)
WARN       = RGBColor(0xB4, 0x53, 0x09)
WARN_BG    = RGBColor(0xFE, 0xF3, 0xC7)
GOOD       = RGBColor(0x16, 0x65, 0x34)
GOOD_BG    = RGBColor(0xDC, 0xFC, 0xE7)
BAD        = RGBColor(0xB9, 0x1C, 0x1C)
BAD_BG     = RGBColor(0xFE, 0xE2, 0xE2)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Margins.
MX = Inches(0.7)
MY_TOP = Inches(0.55)

def new_deck() -> Presentation:
    p = Presentation()
    p.slide_width = SLIDE_W
    p.slide_height = SLIDE_H
    return p

def blank(p: Presentation):
    return p.slides.add_slide(p.slide_layouts[6])  # 6 = blank

def add_textbox(slide, x, y, w, h, text, *, size=14, bold=False, color=INK,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Helvetica"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p_ = tf.paragraphs[0]
    p_.alignment = align
    run = p_.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_multiline(slide, x, y, w, h, lines, *, size=12, color=INK,
                  bold_first=False, font="Helvetica", line_space=1.25,
                  bullet=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.line_spacing = line_space
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = ("• " if bullet else "") + line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        if bold_first and i == 0:
            run.font.bold = True
    return tb

def add_eyebrow(slide, text):
    add_textbox(slide, MX, MY_TOP, Inches(8), Inches(0.3),
                text.upper(), size=11, bold=True, color=ACCENT)

def add_title(slide, text, *, y=None):
    add_textbox(slide, MX, y or Inches(0.9), Inches(12), Inches(0.7),
                text, size=32, bold=True, color=INK)

def add_lede(slide, text, *, y):
    add_textbox(slide, MX, y, Inches(12), Inches(0.8),
                text, size=14, color=INK_SOFT)

def add_section_heading(slide, text, *, x, y, w=Inches(6)):
    add_textbox(slide, x, y, w, Inches(0.3),
                text.upper(), size=10, bold=True, color=INK_SOFT)

def add_pill(slide, x, y, label, kind="live"):
    fill, ink = {
        "live":    (GOOD_BG, GOOD),
        "preview": (WARN_BG, WARN),
        "planned": (BAD_BG,  BAD),
    }[kind]
    w = Inches(0.9)
    h = Inches(0.28)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p_ = tf.paragraphs[0]
    p_.alignment = PP_ALIGN.CENTER
    run = p_.add_run()
    run.text = label.upper()
    run.font.name = "Helvetica"
    run.font.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = ink

def add_card(slide, x, y, w, h, *, title, body):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = RULE
    rect.shadow.inherit = False
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(120000); tf.margin_right = Emu(120000)
    tf.margin_top = Emu(100000); tf.margin_bottom = Emu(100000)
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Helvetica"; r1.font.bold = True
    r1.font.size = Pt(14); r1.font.color.rgb = INK
    p2 = tf.add_paragraph()
    p2.line_spacing = 1.25
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Helvetica"; r2.font.size = Pt(11)
    r2.font.color.rgb = INK_SOFT

def add_callout(slide, x, y, w, h, text, *, kind="accent"):
    fill, bar = {
        "accent": (ACCENT_BG, ACCENT),
        "warn":   (WARN_BG, WARN),
    }[kind]
    bar_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar_shape.fill.solid(); bar_shape.fill.fore_color.rgb = bar
    bar_shape.line.fill.background()
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.08), y,
                                  w - Inches(0.08), h)
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(180000); tf.margin_right = Emu(180000)
    tf.margin_top = Emu(120000); tf.margin_bottom = Emu(120000)
    p_ = tf.paragraphs[0]
    p_.line_spacing = 1.25
    run = p_.add_run()
    run.text = text
    run.font.name = "Helvetica"
    run.font.size = Pt(11)
    run.font.color.rgb = INK

def add_table(slide, x, y, w, h, header, rows, *,
              header_fill=CARD_BG, header_ink=INK_SOFT,
              body_size=10, header_size=9):
    cols = len(header)
    n_rows = 1 + len(rows)
    tbl = slide.shapes.add_table(n_rows, cols, x, y, w, h).table
    for i, label in enumerate(header):
        cell = tbl.cell(0, i)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_left = Emu(100000); cell.margin_right = Emu(100000)
        cell.margin_top = Emu(60000); cell.margin_bottom = Emu(60000)
        tf = cell.text_frame
        p_ = tf.paragraphs[0]
        run = p_.add_run()
        run.text = label.upper()
        run.font.name = "Helvetica"; run.font.bold = True
        run.font.size = Pt(header_size); run.font.color.rgb = header_ink
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = tbl.cell(r_idx, c_idx)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            cell.margin_left = Emu(100000); cell.margin_right = Emu(100000)
            cell.margin_top = Emu(60000); cell.margin_bottom = Emu(60000)
            tf = cell.text_frame
            tf.word_wrap = True
            p_ = tf.paragraphs[0]
            p_.line_spacing = 1.2
            run = p_.add_run()
            run.text = str(value)
            run.font.name = "Helvetica"
            run.font.size = Pt(body_size)
            run.font.color.rgb = INK
    return tbl


# ─────────────────────────────────────────────────────────────────────
# Build the deck
# ─────────────────────────────────────────────────────────────────────

deck = new_deck()


# Slide 1 — Cover
s = blank(deck)
add_eyebrow(s, "BookLets")
add_textbox(s, MX, Inches(1.0), Inches(12), Inches(2.0),
            "A guided tour for the bookkeeper & accountant",
            size=40, bold=True, color=INK)
add_textbox(s, MX, Inches(3.0), Inches(12), Inches(1.6),
            "BookLets is the in-house accounting and operations system for the "
            "Ko Lake short-term-rental portfolio. This deck walks through each "
            "screen, what you'll use it for, and the accounting policies the "
            "system enforces.",
            size=15, color=INK_SOFT)
add_multiline(s, MX, Inches(5.4), Inches(12), Inches(1.4),
              ["Audience:  Bookkeeper, accountant, external reviewer",
               "Currency:  Books in LKR, reporting in USD (month-close FX)",
               "Framework:  SLFRS (Sri Lanka Financial Reporting Standards)"],
              size=12, color=INK_SOFT, line_space=1.4)


# Slide 2 — Contents + legend
s = blank(deck)
add_eyebrow(s, "Contents")
add_title(s, "What's in this pack")
add_table(s, MX, Inches(1.8), Inches(8), Inches(3.5),
          ["Section", "Page"],
          [
              ("1. What BookLets is for", "3"),
              ("2. The monthly workflow at a glance", "4"),
              ("3. Sign-in", "5"),
              ("4. Dashboard", "6"),
              ("5. Properties", "7"),
              ("6. Bookings", "8"),
              ("7. General Ledger", "9"),
              ("8. Imports — spreadsheet uploader", "10"),
              ("9. Chart of accounts & policies", "11"),
              ("10. Roadmap", "12"),
          ],
          body_size=11)
add_section_heading(s, "Status legend", x=Inches(9.4), y=Inches(1.8))
add_pill(s, Inches(9.4), Inches(2.25), "Live", "live")
add_textbox(s, Inches(10.45), Inches(2.27), Inches(3), Inches(0.3),
            "Available in production now.", size=11, color=INK_SOFT)
add_pill(s, Inches(9.4), Inches(2.75), "Preview", "preview")
add_textbox(s, Inches(10.45), Inches(2.77), Inches(3), Inches(0.3),
            "Read-only; posting comes later.", size=11, color=INK_SOFT)
add_pill(s, Inches(9.4), Inches(3.25), "Planned", "planned")
add_textbox(s, Inches(10.45), Inches(3.27), Inches(3), Inches(0.3),
            "Designed, not yet built.", size=11, color=INK_SOFT)


# Slide 3 — What it does
s = blank(deck)
add_eyebrow(s, "Section 1")
add_title(s, "What BookLets is for")
add_lede(s, "A single place to record income and expenses, attribute revenue "
            "to the correct month, produce a clean general ledger, and hand "
            "monthly numbers to the accountant for QuickBooks Online.",
         y=Inches(1.7))

card_w = Inches(5.85); card_h = Inches(1.45); gap = Inches(0.3)
row_y = Inches(3.0)
add_card(s, MX, row_y, card_w, card_h,
         title="Replaces",
         body="Loose spreadsheets for petty-cash analysis; manual re-keying "
              "of monthly totals; ad-hoc reconciliation between operator and "
              "accountant.")
add_card(s, MX + card_w + gap, row_y, card_w, card_h,
         title="Feeds",
         body="QuickBooks Online via CSV export. BookLets is the source of "
              "truth; QBO is the filing surface.")
add_card(s, MX, row_y + card_h + gap, card_w, card_h,
         title="Books in",
         body="LKR (Sri Lankan rupees). Every transaction recorded in LKR "
              "exactly as it happened.")
add_card(s, MX + card_w + gap, row_y + card_h + gap, card_w, card_h,
         title="Reports in",
         body="USD using the spot FX rate on the day the month is closed "
              "(preferred: monthly average if a daily feed is available).")

add_callout(s, MX, Inches(6.45), Inches(12), Inches(0.85),
            "Golden rule — When a guest's stay crosses a month boundary, the "
            "entire stay revenue is booked to the CHECK-OUT month. No "
            "day-by-day apportionment. Refund risk effectively closes on check-out.")


# Slide 4 — Workflow
s = blank(deck)
add_eyebrow(s, "Section 2")
add_title(s, "The monthly workflow")
add_table(s, MX, Inches(1.7), Inches(12), Inches(4.8),
          ["#", "Step", "Where in BookLets", "Owner"],
          [
              ("1", "Record petty-cash and operating expenses on the monthly workbook (Income & Petty Cash Analysis).",
               "External spreadsheet (today)", "Villa captain"),
              ("2", "At month end, upload the workbook. BookLets parses every row and shows a preview.",
               "/imports", "Bookkeeper"),
              ("3", "Review the preview: per-section totals, unmapped columns, rows flagged with warnings.",
               "/imports preview", "Bookkeeper"),
              ("4", "Confirm & post. BookLets writes balanced double-entry journals. Idempotent on re-upload.",
               "/imports (P2 — next)", "Bookkeeper"),
              ("5", "Reconcile to bank. Match cleared transactions; queue exceptions.",
               "/reconcile (P4)", "Bookkeeper"),
              ("6", "Run FX revaluation at month close (LKR → USD).",
               "Month-close (P5)", "Accountant"),
              ("7", "Export ledger CSV and import into QuickBooks Online.",
               "Export CSV button", "Accountant"),
          ],
          body_size=10)


# Slide 5 — Sign-in
s = blank(deck)
add_eyebrow(s, "Screen 1")
add_title(s, "Sign-in")
add_pill(s, Inches(4.2), Inches(1.05), "Live", "live")
add_lede(s, "Access is by Google account only, and only emails on the "
            "allow-list can sign in.", y=Inches(1.7))

add_section_heading(s, "What you'll see", x=MX, y=Inches(2.7))
add_multiline(s, MX, Inches(3.05), Inches(12), Inches(1.2),
              ["A centred card with the BookLets brand, a 'Sign in' heading, and one 'Continue with Google' button.",
               "Clear error messages if your email is not on the allow-list, or if the system is misconfigured."],
              size=12, color=INK_SOFT, line_space=1.3, bullet=True)

add_section_heading(s, "What you need to do", x=MX, y=Inches(4.4))
add_multiline(s, MX, Inches(4.75), Inches(12), Inches(1.3),
              ["1.  Send your Google address to the operator.",
               "2.  They add it to the allow-list.",
               "3.  Visit the site, click 'Continue with Google', and you'll land on the Dashboard."],
              size=12, color=INK_SOFT, line_space=1.3)

add_callout(s, MX, Inches(6.2), Inches(12), Inches(0.95),
            "Security note — If the allow-list is accidentally cleared, "
            "sign-in fails for everyone (fail-closed). Deliberate: the "
            "operator would rather lock themselves out than admit an "
            "unauthorised account into the books.",
            kind="warn")


# Slide 6 — Dashboard
s = blank(deck)
add_eyebrow(s, "Screen 2")
add_title(s, "Dashboard")
add_pill(s, Inches(4.3), Inches(1.05), "Live", "live")
add_lede(s, "First screen after sign-in. One-glance view of how the portfolio "
            "is performing this month and year-to-date.", y=Inches(1.7))

card_w = Inches(2.85); card_h = Inches(1.5)
y0 = Inches(2.7)
add_card(s, MX, y0, card_w, card_h,
         title="Total Revenue",
         body="Month-to-date gross income recognised in the ledger.")
add_card(s, MX + (card_w + Inches(0.2)), y0, card_w, card_h,
         title="Net Income",
         body="Revenue minus operating expenses. Margin in the sub-label.")
add_card(s, MX + (card_w + Inches(0.2)) * 2, y0, card_w, card_h,
         title="ADR / RevPAR",
         body="Average daily rate and revenue per available room — STR yield.")
add_card(s, MX + (card_w + Inches(0.2)) * 3, y0, card_w, card_h,
         title="Portfolio Occupancy",
         body="% of available room-nights actually sold this month.")

add_section_heading(s, "Below the cards", x=MX, y=Inches(4.5))
add_multiline(s, MX, Inches(4.85), Inches(12), Inches(1.5),
              ["Receipt uploader — drop a photo or scan; gets attached to an expense (printed/handwritten flag in P8).",
               "Revenue Trend bar chart — gross revenue and net income side by side for recent months.",
               "Property Yield list — per-villa headline numbers, links into the Properties page."],
              size=12, color=INK_SOFT, line_space=1.3, bullet=True)

add_section_heading(s, "Header buttons", x=MX, y=Inches(6.4))
add_multiline(s, MX, Inches(6.75), Inches(12), Inches(0.6),
              ["Download Report — exports the ledger CSV for the period.",
               "+ Create Entry — jumps to the Ledger to add a manual journal entry."],
              size=12, color=INK_SOFT, line_space=1.3, bullet=True)


# Slide 7 — Properties
s = blank(deck)
add_eyebrow(s, "Screen 3")
add_title(s, "Properties")
add_pill(s, Inches(4.05), Inches(1.05), "Live", "live")
add_lede(s, "One card per villa, with the financial picture for that asset.",
         y=Inches(1.7))

add_table(s, MX, Inches(2.3), Inches(12), Inches(3.5),
          ["Card element", "What it tells you"],
          [
              ("Status pill", "Active, paused, or in setup."),
              ("Total Revenue", "Period-to-date gross income for that villa."),
              ("Net Yield", "Revenue minus directly attributable cost."),
              ("ADR", "Average daily rate over occupied nights."),
              ("RevPAR", "Revenue per available room-night (includes vacant nights)."),
              ("Occupancy bar", "Visual % of nights sold vs available."),
              ("Details →", "Drills into the per-property page (bookings, costs, manager)."),
          ],
          body_size=11)

add_section_heading(s, "Top-right action", x=MX, y=Inches(6.0))
add_textbox(s, MX, Inches(6.35), Inches(12), Inches(0.4),
            "Sync Properties — pulls the latest property list from Hostaway "
            "(channel manager) so BookLets stays aligned with live inventory.",
            size=12, color=INK_SOFT)
add_callout(s, MX, Inches(6.85), Inches(12), Inches(0.5),
            "Currently the portfolio is one villa (Ko Lake). The grid layout "
            "supports any number; new villas can be added without code changes.")


# Slide 8 — Bookings
s = blank(deck)
add_eyebrow(s, "Screen 4")
add_title(s, "Bookings")
add_pill(s, Inches(3.85), Inches(1.05), "Live", "live")
add_lede(s, "A single ordered table of every reservation, regardless of which "
            "channel it came through.", y=Inches(1.7))

add_table(s, MX, Inches(2.5), Inches(12), Inches(3.4),
          ["Column", "What it tells you"],
          [
              ("ID", "Hostaway booking reference, or BookLets' internal short ID if created manually."),
              ("Property", "Which villa."),
              ("Channel", "Booking.com, Airbnb, direct, etc."),
              ("Check In / Check Out", "Stay dates. The check-out date determines which month the revenue is booked to."),
              ("Total", "Gross reservation value, in the channel's billing currency."),
              ("Status", "Coloured badge: Confirmed, Completed, Pending, Cancelled."),
          ],
          body_size=11)

add_callout(s, MX, Inches(6.1), Inches(12), Inches(0.9),
            "Revenue recognition reminder — Status moves to Completed on "
            "check-out. Revenue is recognised at that point and falls into "
            "the check-out month, even if the stay started earlier.")


# Slide 9 — Ledger
s = blank(deck)
add_eyebrow(s, "Screen 5")
add_title(s, "General Ledger")
add_pill(s, Inches(5.0), Inches(1.05), "Live", "live")
add_lede(s, "The double-entry journal. Every income and expense in BookLets "
            "shows up here, in chronological order, with full debit/credit "
            "detail.", y=Inches(1.7))

add_table(s, MX, Inches(2.7), Inches(12), Inches(3.0),
          ["Column", "What it tells you"],
          [
              ("Date", "Transaction date (when it actually happened)."),
              ("Reference", "8-char ID linking every line of a single journal entry — they balance to zero."),
              ("Account", "Chart-of-accounts line item, e.g. 4000 Rent Income or 6200 Electricity."),
              ("Memo", "Description as entered by the bookkeeper or imported from the spreadsheet."),
              ("Debit",  "Green column. Increases assets and expenses; decreases income and liabilities."),
              ("Credit", "Red column. Increases income and liabilities; decreases assets and expenses."),
          ],
          body_size=10)

add_section_heading(s, "Filters & export", x=MX, y=Inches(6.0))
add_multiline(s, MX, Inches(6.35), Inches(12), Inches(0.95),
              ["Period filter — dropdown auto-populated with every month that has activity, plus 'All Time'.",
               "Export CSV — produces a file ready to import into QuickBooks Online. Columns map cleanly to QBO's journal-entry import format."],
              size=11, color=INK_SOFT, line_space=1.3, bullet=True)


# Slide 10 — Imports
s = blank(deck)
add_eyebrow(s, "Screen 6")
add_title(s, "Imports — Spreadsheet Uploader")
add_pill(s, Inches(7.85), Inches(1.05), "Preview", "preview")
add_lede(s, "Where you upload the monthly Income & Petty Cash Analysis "
            "workbook. Currently a read-only preview; posting to the ledger "
            "comes in the next phase (P2).", y=Inches(1.7))

add_section_heading(s, "How it works", x=MX, y=Inches(2.8))
add_multiline(s, MX, Inches(3.15), Inches(12), Inches(1.3),
              ["1.  Drop the .xlsx for the month into the uploader. Limit 10 MB.",
               "2.  BookLets parses every row, reads column headers, and matches each amount to a chart-of-accounts code.",
               "3.  Preview appears below the form. Nothing is written to the ledger until the confirm-and-post step."],
              size=11, color=INK_SOFT, line_space=1.3)

add_section_heading(s, "What the preview shows you", x=MX, y=Inches(4.85))
add_table(s, MX, Inches(5.2), Inches(12), Inches(1.9),
          ["Section", "Content"],
          [
              ("Summary card", "Period, total rows, net amount, file fingerprint."),
              ("Totals by account, per section", "Mini-table per section: Prior-Month, Daily, Recurring, Accruals, Reversals, Prepayments."),
              ("Unmapped columns", "Amber banner listing any header BookLets couldn't match to a chart account."),
              ("Per-section row tables", "Every row with date, description, petty cash, postings, warnings."),
          ],
          body_size=10)


# Slide 11 — Chart & policies
s = blank(deck)
add_eyebrow(s, "Section 9")
add_title(s, "Chart of accounts & policies")
add_section_heading(s, "Chart of accounts (36 lines)", x=MX, y=Inches(1.7))
add_table(s, MX, Inches(2.05), Inches(12), Inches(2.4),
          ["Range", "Type", "Examples"],
          [
              ("1xxx", "Assets",        "1010 Petty Cash, 1100 Bank — LKR Current"),
              ("2xxx", "Liabilities",   "2200 APIT Payable, 2210 EPF Payable, 2220 ETF Payable"),
              ("4xxx", "Revenue",       "4000 Rent, 4010 Cleaning Fee, 4020 Event, 4030 F&B"),
              ("5xxx", "Cost of sales", "5100 Food & Beverage Expense, 5110 Refunds"),
              ("6xxx", "Operating",     "6100 Salaries, 6200 Electricity, 6300 Cleaning, 6600 Admin"),
              ("7xxx", "Capex",         "7100 Minor Capex"),
              ("9999", "Suspense",      "Used for unmapped amounts until the operator classifies them."),
          ],
          body_size=10)

add_section_heading(s, "Key accounting policies", x=MX, y=Inches(4.6))
add_table(s, MX, Inches(4.95), Inches(12), Inches(2.2),
          ["Topic", "Policy"],
          [
              ("Currency", "Books in LKR. USD reporting via month-close FX rate (preferred: monthly average)."),
              ("Revenue recognition", "Entire stay revenue booked to the check-out month. No day-by-day apportionment."),
              ("Petty cash", "Held by the villa captain. Top-ups against 1010. Items > LKR 5,000 require a memo."),
              ("Payroll", "Gross to 6100 Salaries; statutory 2200 APIT, 2210 EPF (8% emp + 12% emp'r), 2220 ETF (3% emp'r). Net to 6110."),
              ("Framework", "SLFRS (Sri Lanka Financial Reporting Standards). Swappable if portfolio expands overseas."),
          ],
          body_size=10)


# Slide 12 — Roadmap
s = blank(deck)
add_eyebrow(s, "Section 10")
add_title(s, "Roadmap — what's coming next")
add_lede(s, "In rough priority order. Everything below P1 is designed; "
            "nothing else is in production yet.", y=Inches(1.7))

add_table(s, MX, Inches(2.5), Inches(12), Inches(4.0),
          ["Phase", "Feature", "Status"],
          [
              ("P0",     "Chart of accounts — 36 lines drafted from the operator's workbook",              "Live"),
              ("P1",     "Spreadsheet parser + read-only preview",                                          "Preview"),
              ("P2",     "Confirm-and-post — writes balanced journal entries; idempotent on re-upload",     "Planned"),
              ("P3",     "Editable grid — fix typos and re-classify before posting",                        "Planned"),
              ("P4",     "Bank reconciliation — match ledger entries to cleared bank lines",                "Planned"),
              ("P5",     "Month close — FX revaluation, accountant export pack",                            "Planned"),
              ("P6",     "STR dashboards — ADR, RevPAR, channel mix, seasonality",                          "Planned"),
              ("P7",     "Capex tracker & forecast editor",                                                 "Planned"),
              ("P8",     "Google Drive receipts pipeline — auto-OCR, P/H flag, Sinhala translation",        "Planned"),
              ("P9–11",  "AI chat dialog per screen — database-grounded, SLFRS-sourced for accounting Q&A", "Planned"),
          ],
          body_size=9)

add_callout(s, MX, Inches(6.7), Inches(12), Inches(0.55),
            "Your input is welcomed. If anything looks wrong — a missing "
            "account, a policy you'd state differently, a column heading the "
            "parser should recognise — please flag it.")


# Save
out = Path(__file__).resolve().parent.parent / "docs" / "booklets-walkthrough.pptx"
out.parent.mkdir(exist_ok=True)
deck.save(out)
print(f"wrote {out}")
